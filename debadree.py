import os
import time
import re
import traceback
import speech_recognition as sr
from dotenv import load_dotenv
from langgraph.graph import StateGraph, END
from typing import TypedDict
from pptx import Presentation
import subprocess
import pyttsx3
import google.generativeai as genai

# Load environment variables
load_dotenv()

# Configure Gemini API
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

# Initialize LLM
llm = genai.GenerativeModel('gemini-2.0-flash-exp')


class ConversationState(TypedDict):
    """State for conversation flow"""
    user_input: str
    intent: str
    response: str
    action: str
    slide_number: int


class AIPresenter:
    def _init_(self):
        self.recognizer = sr.Recognizer()
        self.microphone = None
        self.presentation = None
        self.current_slide = 0
        self.total_slides = 0
        self.ppt_file = None
        self.presentation_mode = False
        
        # Try to initialize microphone
        try:
            self.microphone = sr.Microphone()
            print("[INFO] Microphone initialized successfully")
        except Exception as e:
            print(f"[WARNING] Could not initialize microphone: {e}")
            print("[INFO] Voice input will be disabled. You can still use text input.")
        
    def speak(self, text):
        """Convert text to speech with engine reinitialization"""
        print(f"\nAssistant: {text}")
        try:
            # Clean up old engine
            if hasattr(self, 'engine') and self.engine is not None:
                try:
                    self.engine.stop()
                except:
                    pass
            
            # Create fresh engine for each speech
            engine = pyttsx3.init('espeak')
            engine.setProperty('rate', 150)
            engine.setProperty('volume', 0.9)
            engine.say(text)
            engine.runAndWait()
            
            # Clean up
            del engine
            
        except Exception as e:
            print(f"[TTS Error: {e}]")
    
    def listen(self):
        """Listen to user's voice input"""
        if not self.microphone:
            print("[ERROR] Microphone not available. Please use text input.")
            return ""
            
        with self.microphone as source:
            print("\n🎤 Listening...")
            self.recognizer.adjust_for_ambient_noise(source, duration=0.5)
            try:
                audio = self.recognizer.listen(source, timeout=5, phrase_time_limit=10)
                text = self.recognizer.recognize_google(audio)
                print(f"\n👤 You: {text}")
                return text
            except sr.WaitTimeoutError:
                print("[INFO] No speech detected")
                return ""
            except sr.UnknownValueError:
                print("[ERROR] Could not understand audio")
                return ""
            except sr.RequestError as e:
                print(f"[ERROR] Speech recognition service error: {e}")
                return ""
            except Exception as e:
                print(f"[ERROR] Unexpected error: {e}")
                return ""
    
    def load_presentation(self, file_path):
        """Load a PowerPoint presentation"""
        try:
            # Expand user path and resolve relative paths
            file_path = os.path.expanduser(file_path)
            if not os.path.isabs(file_path):
                file_path = os.path.abspath(file_path)
            
            if not os.path.exists(file_path):
                print(f"[ERROR] File not found: {file_path}")
                return False
            
            self.presentation = Presentation(file_path)
            self.total_slides = len(self.presentation.slides)
            self.current_slide = 0
            self.ppt_file = file_path
            self.presentation_mode = True
            print(f"[SUCCESS] Loaded presentation with {self.total_slides} slides")
            return True
        except Exception as e:
            print(f"[ERROR] Failed to load presentation: {str(e)}")
            traceback.print_exc()
            return False
    
    def extract_slide_content(self, slide_number):
        """Extract text content from a specific slide"""
        if not self.presentation or slide_number < 0 or slide_number >= self.total_slides:
            return None
        
        slide = self.presentation.slides[slide_number]
        content = {
            "slide_number": slide_number + 1,
            "title": "",
            "bullets": [],
            "notes": ""
        }
        
        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    # Check if it's a title placeholder
                    if hasattr(shape, "placeholder_format"):
                        if shape.placeholder_format.type == 1:  # Title
                            content["title"] = text
                            continue
                    
                    # Otherwise treat as bullet point
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            bullet_text = paragraph.text.strip()
                            if bullet_text and bullet_text not in content["bullets"]:
                                content["bullets"].append(bullet_text)
        
        # Extract notes
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    content["notes"] = notes_slide.notes_text_frame.text.strip()
        except Exception as e:
            print(f"[WARNING] Could not extract notes: {e}")
        
        return content
    
    def open_presentation(self, slide_number=0):
        """Open the PowerPoint file (on Raspberry Pi, uses LibreOffice)"""
        if not self.ppt_file:
            return False
        
        try:
            abs_path = os.path.abspath(self.ppt_file)
            
            # Try LibreOffice Impress (common on Raspberry Pi)
            try:
                subprocess.Popen(['libreoffice', '--impress', abs_path])
                print("[INFO] Opening presentation with LibreOffice Impress")
                return True
            except FileNotFoundError:
                pass
            
            # Fallback to xdg-open
            try:
                subprocess.Popen(['xdg-open', abs_path])
                print("[INFO] Opening presentation with default application")
                return True
            except:
                pass
            
            print("[WARNING] Could not open presentation file")
            return False
            
        except Exception as e:
            print(f"[ERROR] Failed to open presentation: {str(e)}")
            return False
    
    def classify_intent(self, state: ConversationState) -> ConversationState:
        """Classify user intent as generic chat or presentation command"""
        user_input = state["user_input"].lower()
        
        # Keywords for presentation commands
        presentation_keywords = [
            "slide", "next", "previous", "back", "forward", "go to", "goto",
            "explain", "present", "show", "open presentation", "start presentation",
            "first", "last", "current"
        ]
        
        if self.presentation_mode and any(keyword in user_input for keyword in presentation_keywords):
            state["intent"] = "presentation"
        else:
            state["intent"] = "generic"
        
        return state
    
    def handle_generic_chat(self, state: ConversationState) -> ConversationState:
        """Handle generic chatbot conversation"""
        if state["intent"] != "generic":
            return state
        
        try:
            # Add context if in presentation mode
            context = ""
            if self.presentation_mode:
                context = f"\nNote: We are currently in presentation mode. The presentation has {self.total_slides} slides and we're on slide {self.current_slide + 1}. "
            
            response = llm.generate_content(context + state["user_input"])
            state["response"] = response.text
            state["action"] = "none"
        except Exception as e:
            state["response"] = f"I encountered an error: {str(e)}"
            state["action"] = "none"
            print(f"[ERROR] LLM error: {e}")
        
        return state
    
    def parse_presentation_command(self, state: ConversationState) -> ConversationState:
        """Parse presentation navigation commands"""
        if state["intent"] != "presentation":
            return state
        
        user_input = state["user_input"].lower()
        
        # Extract slide number if mentioned
        numbers = re.findall(r'\d+', user_input)
        target_slide = None
        
        if numbers:
            target_slide = int(numbers[0]) - 1  # Convert to 0-indexed
        
        # Determine action
        if "next" in user_input:
            state["action"] = "next"
            state["slide_number"] = self.current_slide + 1
        elif "previous" in user_input or "back" in user_input:
            state["action"] = "previous"
            state["slide_number"] = self.current_slide - 1
        elif "first" in user_input:
            state["action"] = "goto"
            state["slide_number"] = 0
        elif "last" in user_input:
            state["action"] = "goto"
            state["slide_number"] = self.total_slides - 1
        elif "current" in user_input:
            state["action"] = "explain"
            state["slide_number"] = self.current_slide
        elif ("go to" in user_input or "goto" in user_input) and target_slide is not None:
            state["action"] = "goto"
            state["slide_number"] = target_slide
        elif "explain" in user_input and target_slide is not None:
            state["action"] = "explain"
            state["slide_number"] = target_slide
        elif "explain" in user_input:
            state["action"] = "explain"
            state["slide_number"] = self.current_slide
        elif "open" in user_input or "start" in user_input:
            state["action"] = "open"
            state["slide_number"] = self.current_slide
        elif target_slide is not None:
            # Just a number mentioned, go to that slide
            state["action"] = "goto"
            state["slide_number"] = target_slide
        else:
            state["action"] = "unknown"
            state["slide_number"] = self.current_slide
        
        return state
    
    def execute_presentation_action(self, state: ConversationState) -> ConversationState:
        """Execute the presentation action"""
        if state["intent"] != "presentation":
            return state
        
        action = state["action"]
        slide_num = state["slide_number"]
        
        # Validate slide number
        if slide_num < 0:
            state["response"] = "We're already at the first slide."
            return state
        
        if slide_num >= self.total_slides:
            state["response"] = f"There are only {self.total_slides} slides in this presentation."
            return state
        
        # Execute action
        if action == "open":
            if self.open_presentation(slide_num):
                state["response"] = f"Opening the presentation. We have {self.total_slides} slides."
            else:
                state["response"] = "Sorry, I couldn't open the presentation file."
            return state
        
        if action == "unknown":
            state["response"] = "I didn't understand that command. You can say things like 'next slide', 'go to slide 5', or 'explain this slide'."
            return state
        
        # Navigate to slide
        self.current_slide = slide_num
        slide_content = self.extract_slide_content(slide_num)
        
        if not slide_content:
            state["response"] = "Sorry, I couldn't read that slide."
            return state
        
        # Generate response based on action
        if action in ["next", "previous", "goto"]:
            response = f"Now on slide {slide_content['slide_number']} of {self.total_slides}. "
            if slide_content['title']:
                response += f"Title: {slide_content['title']}. "
            
            if action != "goto" or "explain" not in state["user_input"].lower():
                # Just announce the slide
                state["response"] = response
            else:
                # Also explain it
                state["response"] = response + self.explain_slide(slide_content)
        
        elif action == "explain":
            state["response"] = self.explain_slide(slide_content)
        
        return state
    
    def explain_slide(self, slide_content):
        """Generate an explanation for the slide content using LLM"""
        try:
            prompt = f"""You are presenting a PowerPoint slide. Here is the content:

Slide {slide_content['slide_number']}:
Title: {slide_content['title']}

Bullet Points:
{chr(10).join(['- ' + bullet for bullet in slide_content['bullets']])}

{"Speaker Notes: " + slide_content['notes'] if slide_content['notes'] else ""}

Provide a clear, engaging explanation of this slide as if you're presenting it to an audience. Be concise but informative."""

            response = llm.generate_content(prompt)
            return response.text
        
        except Exception as e:
            print(f"[ERROR] Failed to generate explanation: {e}")
            # Fallback explanation
            explanation = f"This is slide {slide_content['slide_number']}. "
            if slide_content['title']:
                explanation += f"The title is '{slide_content['title']}'. "
            
            if slide_content['bullets']:
                explanation += "The main points are: " + ", ".join(slide_content['bullets'][:3])
                if len(slide_content['bullets']) > 3:
                    explanation += f", and {len(slide_content['bullets']) - 3} more points."
            
            return explanation
    
    def route_intent(self, state: ConversationState) -> str:
        """Route to appropriate handler based on intent"""
        if state["intent"] == "presentation":
            return "presentation"
        else:
            return "generic"
    
    def build_graph(self):
        """Build the conversation flow graph"""
        workflow = StateGraph(ConversationState)
        
        # Add nodes
        workflow.add_node("classify", self.classify_intent)
        workflow.add_node("generic", self.handle_generic_chat)
        workflow.add_node("parse_command", self.parse_presentation_command)
        workflow.add_node("execute_action", self.execute_presentation_action)
        
        # Add edges
        workflow.set_entry_point("classify")
        workflow.add_conditional_edges(
            "classify",
            self.route_intent,
            {
                "generic": "generic",
                "presentation": "parse_command"
            }
        )
        workflow.add_edge("generic", END)
        workflow.add_edge("parse_command", "execute_action")
        workflow.add_edge("execute_action", END)
        
        return workflow.compile()
    
    def run(self):
        """Main run loop"""
        graph = self.build_graph()
        
        self.speak("Hello! I'm your AI presenter assistant. I can chat with you or help present PowerPoint slides.")
        
        # Load presentation
        ppt_path = input("\n📁 Enter the path to your PowerPoint file (or press Enter to skip): ").strip()
        
        if ppt_path:
            if self.load_presentation(ppt_path):
                self.speak(f"Presentation loaded successfully. It has {self.total_slides} slides. Say 'open presentation' to start, or ask me to navigate.")
            else:
                self.speak("Failed to load the presentation. You can still chat with me.")
        else:
            self.speak("No presentation loaded. You can chat with me or load a presentation later.")
        
        while True:
            # Get user input (voice or text)
            print("\n" + "="*60)
            
            input_options = "Choose input: "
            if self.microphone:
                input_options += "1=Voice, 2=Text, 3=Exit"
            else:
                input_options += "2=Text, 3=Exit (Voice disabled)"
            
            choice = input(input_options + ": ").strip()
            
            if choice == "3":
                self.speak("Goodbye! Have a great day!")
                break
            
            if choice == "1":
                if not self.microphone:
                    print("[ERROR] Voice input not available. Please use text input (option 2).")
                    continue
                user_input = self.listen()
                if not user_input:
                    continue
            elif choice == "2":
                user_input = input("\n👤 You: ").strip()
                if not user_input:
                    continue
            else:
                print("[ERROR] Invalid choice. Please try again.")
                continue
            
            # Check for exit commands
            if user_input.lower() in ["exit", "quit", "bye", "goodbye"]:
                self.speak("Goodbye! Have a great day!")
                break
            
            # Check for load presentation command
            if "load presentation" in user_input.lower():
                ppt_path = input("📁 Enter the path to your PowerPoint file: ").strip()
                if self.load_presentation(ppt_path):
                    self.speak(f"Presentation loaded successfully. It has {self.total_slides} slides.")
                else:
                    self.speak("Failed to load the presentation.")
                continue
            
            # Process through graph
            state = {
                "user_input": user_input,
                "intent": "",
                "response": "",
                "action": "none",
                "slide_number": self.current_slide
            }
            
            try:
                result = graph.invoke(state)
                # Respond
                self.speak(result["response"])
            except Exception as e:
                print(f"[ERROR] Processing error: {e}")
                traceback.print_exc()
                self.speak("Sorry, I encountered an error processing your request.")


def main():
    """Main entry point"""
    print("="*60)
    print("🤖 AI Presenter - Chatbot & PPT Navigator")
    print("   Optimized for Raspberry Pi 4")
    print("="*60)
    
    # Check for API key
    if not os.getenv("GEMINI_API_KEY"):
        print("\n❌ Error: GEMINI_API_KEY not found in .env file")
        print("Please create a .env file with:")
        print("GEMINI_API_KEY=your_api_key_here")
        return
    
    presenter = AIPresenter()
    
    try:
        presenter.run()
    except KeyboardInterrupt:
        print("\n\n⚠  Program interrupted by user.")
    except Exception as e:
        print(f"\n❌ An error occurred: {str(e)}")
        traceback.print_exc()


if __name__ == "__main__":
    main()