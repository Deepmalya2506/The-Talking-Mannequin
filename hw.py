"""
Standalone Laptop Version - Voice Agent + PowerPoint Control + RAG
Runs entirely on Windows laptop with document Q&A capabilities
"""

import os
import time
import queue
import threading
import traceback
import speech_recognition as sr
import pyttsx3
from dotenv import load_dotenv
from langchain_groq import ChatGroq
from langgraph.graph import StateGraph, END
from typing import TypedDict, List, Dict, Optional
from datetime import datetime
from zoneinfo import ZoneInfo
import win32com.client
from pathlib import Path

# RAG imports
try:
    import PyPDF2
    from pptx import Presentation as PptxPresentation
    import chromadb
    from sentence_transformers import SentenceTransformer
    RAG_AVAILABLE = True
except ImportError as e:
    print(f"⚠️ RAG dependencies not available: {e}")
    print("Install: pip install PyPDF2 python-pptx chromadb sentence-transformers")
    RAG_AVAILABLE = False

# ---------------- CONFIG ----------------
load_dotenv()
DOCUMENTS_FOLDER = Path(r"C:\Users\DEEPMALYA\OneDrive\Desktop\pip_Malya")
VECTOR_DB_PATH = "./chroma_db"
TIMEZONE = 'Asia/Kolkata'
WAKE_WORDS = ['hello', 'wake up', 'hey', 'jarvis']

# Create vector DB folder
Path(VECTOR_DB_PATH).mkdir(exist_ok=True)

# Presentation state
current_document = None
current_page = 0
ppt_app = None
presentation = None
slide_show = None

# ---------------- SHARED MEMORY ----------------
conversation_history = []
user_input_queue = queue.Queue()
shutdown_flag = threading.Event()
is_agent_awake = threading.Event()
has_greeted = threading.Event()
is_speaking = threading.Event()
speech_interrupted = threading.Event()

# ---------------- LLM MODEL ----------------
llm = ChatGroq(model="llama-3.3-70b-versatile", temperature=0.7)

def model_call(prompt: str) -> str:
    """Centralized LLM call with error handling"""
    try:
        response = llm.invoke(prompt)
        return response.content.strip()
    except Exception as e:
        print(f"[Model Error] {e}")
        return "I'm having trouble processing that right now."

# ---------------- RAG DOCUMENT MANAGER CLASS ----------------
class RAGDocumentManager:
    """
    Manages document indexing and retrieval for intelligent Q&A
    
    Features:
    - Ingests PDF and PowerPoint files
    - Slices documents into chunks/pages
    - Creates embeddings using SentenceTransformers
    - Stores in ChromaDB vector database
    - Retrieves relevant context for queries
    """
    
    def __init__(self, vector_db_path: str = VECTOR_DB_PATH):
        self.vector_db_path = vector_db_path
        self.chroma_client = None
        self.collection = None
        self.embedding_model = None
        self.indexed_documents = {}  # {filename: num_chunks}
        self.document_metadata = {}  # {filename: {total_pages, file_type, indexed_date}}
        
        if RAG_AVAILABLE:
            self.initialize()
    
    def initialize(self) -> bool:
        """Initialize ChromaDB and embedding model"""
        try:
            print("🔄 Initializing RAG system...")
            
            # Initialize ChromaDB with persistent storage
            self.chroma_client = chromadb.PersistentClient(path=self.vector_db_path)
            
            # Create or get collection
            self.collection = self.chroma_client.get_or_create_collection(
                name="document_qa",
                metadata={"hnsw:space": "cosine"}  # Cosine similarity for semantic search
            )
            
            # Load lightweight embedding model (runs offline)
            print("📥 Loading embedding model (all-MiniLM-L6-v2)...")
            self.embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
            
            print("✅ RAG system initialized successfully")
            return True
            
        except Exception as e:
            print(f"❌ RAG initialization failed: {e}")
            traceback.print_exc()
            return False
    
    # ========== DOCUMENT INGESTION ==========
    
    def ingest_document(self, filename: str) -> Dict:
        """
        Main ingestion pipeline: load → extract → chunk → embed → store
        
        Args:
            filename: Name of file in DOCUMENTS_FOLDER
            
        Returns:
            Dict with success status and metadata
        """
        if not RAG_AVAILABLE or not self.embedding_model:
            return {"success": False, "error": "RAG not available"}
        
        try:
            filepath = DOCUMENTS_FOLDER / filename
            
            if not filepath.exists():
                return {"success": False, "error": f"File not found: {filename}"}
            
            # Check if already indexed
            if filename in self.indexed_documents:
                return {
                    "success": True, 
                    "message": "Already indexed",
                    "chunks": self.indexed_documents[filename],
                    "metadata": self.document_metadata.get(filename, {})
                }
            
            print(f"\n📄 Ingesting document: {filename}")
            
            # Step 1: Extract text based on file type
            file_ext = filepath.suffix.lower()
            if file_ext in ['.ppt', '.pptx']:
                pages_data = self._extract_from_pptx(filepath)
            elif file_ext == '.pdf':
                pages_data = self._extract_from_pdf(filepath)
            else:
                return {"success": False, "error": f"Unsupported file type: {file_ext}"}
            
            if not pages_data:
                return {"success": False, "error": "No text extracted"}
            
            # Step 2: Slice into chunks and create embeddings
            result = self._index_chunks(filename, pages_data)
            
            if result['success']:
                # Store metadata
                self.document_metadata[filename] = {
                    "total_pages": len(pages_data),
                    "file_type": file_ext,
                    "indexed_date": datetime.now(ZoneInfo(TIMEZONE)).isoformat(),
                    "chunks": result['chunks']
                }
            
            return result
            
        except Exception as e:
            print(f"❌ Error ingesting document: {e}")
            traceback.print_exc()
            return {"success": False, "error": str(e)}
    
    def _extract_from_pptx(self, filepath: Path) -> List[Dict]:
        """Extract text from PowerPoint slides"""
        pages_data = []
        try:
            prs = PptxPresentation(str(filepath))
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts = []
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text.strip())
                
                # Only add slides with content
                if text_parts:
                    pages_data.append({
                        "page_num": slide_num,
                        "text": "\n".join(text_parts),
                        "type": "slide"
                    })
            
            print(f"  ✅ Extracted {len(pages_data)} slides")
            return pages_data
            
        except Exception as e:
            print(f"  ❌ PowerPoint extraction error: {e}")
            return []
    
    def _extract_from_pdf(self, filepath: Path) -> List[Dict]:
        """Extract text from PDF pages"""
        pages_data = []
        try:
            with open(filepath, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    text = page.extract_text()
                    
                    # Only add pages with meaningful content
                    if text.strip():
                        pages_data.append({
                            "page_num": page_num,
                            "text": text.strip(),
                            "type": "page"
                        })
            
            print(f"  ✅ Extracted {len(pages_data)} pages")
            return pages_data
            
        except Exception as e:
            print(f"  ❌ PDF extraction error: {e}")
            return []
    
    # ========== INDEXING & EMBEDDING ==========
    
    def _index_chunks(self, filename: str, pages_data: List[Dict]) -> Dict:
        """
        Create embeddings and store in ChromaDB
        
        Args:
            filename: Document filename
            pages_data: List of extracted page/slide data
            
        Returns:
            Dict with indexing results
        """
        try:
            print(f"  🔄 Creating embeddings for {len(pages_data)} chunks...")
            
            documents = []
            metadatas = []
            ids = []
            
            # Prepare data for embedding
            for page_data in pages_data:
                text = page_data['text']
                
                # Skip empty chunks
                if not text.strip():
                    continue
                
                # Optionally split long pages into smaller chunks
                chunks = self._split_into_chunks(text, max_length=500)
                
                for chunk_idx, chunk in enumerate(chunks):
                    documents.append(chunk)
                    metadatas.append({
                        "filename": filename,
                        "page_num": str(page_data['page_num']),
                        "chunk_idx": str(chunk_idx),
                        "type": page_data['type']
                    })
                    ids.append(f"{filename}_p{page_data['page_num']}_c{chunk_idx}")
            
            if not documents:
                return {"success": False, "error": "No valid chunks to index"}
            
            # Generate embeddings (this is where the magic happens!)
            print(f"  🧠 Generating embeddings...")
            embeddings = self.embedding_model.encode(
                documents,
                show_progress_bar=False,
                convert_to_numpy=True
            ).tolist()
            
            # Store in ChromaDB
            print(f"  💾 Storing in vector database...")
            self.collection.add(
                embeddings=embeddings,
                documents=documents,
                metadatas=metadatas,
                ids=ids
            )
            
            # Track indexed document
            self.indexed_documents[filename] = len(documents)
            
            print(f"  ✅ Successfully indexed {len(documents)} chunks")
            return {
                "success": True,
                "chunks": len(documents),
                "pages": len(pages_data),
                "filename": filename
            }
            
        except Exception as e:
            print(f"  ❌ Indexing error: {e}")
            traceback.print_exc()
            return {"success": False, "error": str(e)}
    
    def _split_into_chunks(self, text: str, max_length: int = 500) -> List[str]:
        """Split long text into smaller chunks for better retrieval"""
        # Simple sentence-based chunking
        sentences = text.replace('\n', ' ').split('. ')
        
        chunks = []
        current_chunk = []
        current_length = 0
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
            
            sentence_length = len(sentence)
            
            if current_length + sentence_length > max_length and current_chunk:
                # Save current chunk and start new one
                chunks.append('. '.join(current_chunk) + '.')
                current_chunk = [sentence]
                current_length = sentence_length
            else:
                current_chunk.append(sentence)
                current_length += sentence_length
        
        # Add remaining chunk
        if current_chunk:
            chunks.append('. '.join(current_chunk) + '.')
        
        return chunks if chunks else [text]  # Return original if splitting fails
    
    # ========== RETRIEVAL ==========
    
    def retrieve_context(
        self, 
        query: str, 
        filename: Optional[str] = None, 
        top_k: int = 3
    ) -> List[Dict]:
        """
        Retrieve relevant context for a query using semantic search
        
        Args:
            query: User's question
            filename: Optional filter for specific document
            top_k: Number of chunks to retrieve
            
        Returns:
            List of relevant chunks with metadata
        """
        if not RAG_AVAILABLE or not self.embedding_model:
            print("⚠️ RAG not available for retrieval")
            return []
        
        try:
            print(f"  🔍 Searching for: '{query}'")
            
            # Generate query embedding
            query_embedding = self.embedding_model.encode([query])[0].tolist()
            
            # Build filter for specific document
            where_filter = {"filename": filename} if filename else None
            
            # Query ChromaDB (semantic search!)
            results = self.collection.query(
                query_embeddings=[query_embedding],
                n_results=top_k,
                where=where_filter
            )
            
            # Format results
            formatted_results = []
            if results['documents'] and results['documents'][0]:
                for i, doc in enumerate(results['documents'][0]):
                    formatted_results.append({
                        "text": doc,
                        "page_num": results['metadatas'][0][i]['page_num'],
                        "filename": results['metadatas'][0][i]['filename'],
                        "relevance_score": 1 - results['distances'][0][i] if 'distances' in results else None
                    })
                
                print(f"  ✅ Found {len(formatted_results)} relevant chunks")
            else:
                print(f"  ⚠️ No relevant chunks found")
            
            return formatted_results
            
        except Exception as e:
            print(f"  ❌ Retrieval error: {e}")
            traceback.print_exc()
            return []
    
    # ========== HELPER METHODS ==========
    
    def get_indexed_documents(self) -> List[str]:
        """Get list of all indexed documents"""
        return list(self.indexed_documents.keys())
    
    def get_document_info(self, filename: str) -> Optional[Dict]:
        """Get metadata about an indexed document"""
        return self.document_metadata.get(filename)
    
    def is_indexed(self, filename: str) -> bool:
        """Check if document is already indexed"""
        return filename in self.indexed_documents
    
    def clear_index(self, filename: Optional[str] = None):
        """Clear index for specific document or all documents"""
        try:
            if filename:
                # Clear specific document
                if filename in self.indexed_documents:
                    # Delete from ChromaDB
                    ids_to_delete = [
                        doc_id for doc_id in self.collection.get()['ids']
                        if doc_id.startswith(filename)
                    ]
                    if ids_to_delete:
                        self.collection.delete(ids=ids_to_delete)
                    
                    del self.indexed_documents[filename]
                    if filename in self.document_metadata:
                        del self.document_metadata[filename]
                    
                    print(f"✅ Cleared index for {filename}")
            else:
                # Clear all
                self.collection.delete(where={})
                self.indexed_documents.clear()
                self.document_metadata.clear()
                print("✅ Cleared all indexes")
                
        except Exception as e:
            print(f"❌ Error clearing index: {e}")

# Initialize RAG manager
rag_manager = RAGDocumentManager() if RAG_AVAILABLE else None

# ---------------- POWERPOINT CONTROLLER ----------------
class PowerPointController:
    """Direct PowerPoint control on same machine"""
    
    def __init__(self):
        self.ppt_app = None
        self.presentation = None
        self.slide_show = None
        self.initialize()
    
    def initialize(self):
        """Initialize PowerPoint application"""
        try:
            self.ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            self.ppt_app.Visible = True
            print("✅ PowerPoint initialized")
            return True
        except Exception as e:
            print(f"❌ Failed to initialize PowerPoint: {e}")
            return False
    
    def open_presentation(self, filename: str) -> bool:
        """Open a presentation file in full screen slideshow"""
        try:
            filepath = DOCUMENTS_FOLDER / filename
            
            if not filepath.exists():
                print(f"❌ File not found: {filepath}")
                return False
            
            # Close existing presentation
            if self.presentation:
                try:
                    if self.slide_show:
                        self.slide_show.View.Exit()
                        self.slide_show = None
                    self.presentation.Close()
                except:
                    pass
            
            # Open new presentation
            self.presentation = self.ppt_app.Presentations.Open(
                str(filepath.absolute()),
                ReadOnly=False,
                Untitled=False,
                WithWindow=True
            )
            
            # Bring PowerPoint to foreground
            self.ppt_app.Activate()
            
            # Configure slideshow settings for foreground display
            settings = self.presentation.SlideShowSettings
            settings.ShowType = 1
            settings.ShowWithNarration = False
            settings.ShowWithAnimation = True
            settings.AdvanceMode = 0
            
            # Start slideshow and bring to front
            settings.Run()
            
            # Get slideshow window and force to foreground
            self.slide_show = self.presentation.SlideShowWindow
            
            # Force window to foreground using Windows API
            import ctypes
            try:
                hwnd = self.slide_show.HWND
                SW_RESTORE = 9
                ctypes.windll.user32.ShowWindow(hwnd, SW_RESTORE)
                ctypes.windll.user32.SetForegroundWindow(hwnd)
                ctypes.windll.user32.BringWindowToTop(hwnd)
                ctypes.windll.user32.SetFocus(hwnd)
            except Exception as e:
                print(f"⚠️ Could not force to foreground: {e}")
            
            print(f"✅ Opened: {filename} ({self.presentation.Slides.Count} slides)")
            print(f"🖥️ Presentation displayed in foreground")
            return True
            
        except Exception as e:
            print(f"❌ Error opening presentation: {e}")
            traceback.print_exc()
            return False
    
    def next_slide(self) -> bool:
        """Go to next slide and ensure window stays in foreground"""
        try:
            if not self.slide_show:
                return False
            
            view = self.slide_show.View
            if view.Slide.SlideIndex < self.presentation.Slides.Count:
                view.Next()
                self._bring_to_front()
                print(f"➡️ Slide {view.Slide.SlideIndex}")
                return True
            else:
                print("⚠️ Last slide")
                return False
        except Exception as e:
            print(f"❌ Error: {e}")
            return False
    
    def previous_slide(self) -> bool:
        """Go to previous slide and ensure window stays in foreground"""
        try:
            if not self.slide_show:
                return False
            
            view = self.slide_show.View
            if view.Slide.SlideIndex > 1:
                view.Previous()
                self._bring_to_front()
                print(f"⬅️ Slide {view.Slide.SlideIndex}")
                return True
            else:
                print("⚠️ First slide")
                return False
        except Exception as e:
            print(f"❌ Error: {e}")
            return False
    
    def goto_slide(self, slide_number: int) -> bool:
        """Go to specific slide and ensure window stays in foreground"""
        try:
            if not self.slide_show:
                return False
            
            total = self.presentation.Slides.Count
            if 1 <= slide_number <= total:
                self.slide_show.View.GotoSlide(slide_number)
                self._bring_to_front()
                print(f"🎯 Slide {slide_number}")
                return True
            else:
                print(f"⚠️ Invalid slide: {slide_number}")
                return False
        except Exception as e:
            print(f"❌ Error: {e}")
            return False
    
    def _bring_to_front(self):
        """Helper method to bring slideshow window to foreground"""
        try:
            if self.slide_show:
                import ctypes
                hwnd = self.slide_show.HWND
                SW_RESTORE = 9
                ctypes.windll.user32.ShowWindow(hwnd, SW_RESTORE)
                ctypes.windll.user32.SetForegroundWindow(hwnd)
                ctypes.windll.user32.BringWindowToTop(hwnd)
        except Exception as e:
            pass
    
    def close_presentation(self) -> bool:
        """Close presentation"""
        try:
            if self.slide_show:
                self.slide_show.View.Exit()
                self.slide_show = None
            
            if self.presentation:
                self.presentation.Close()
                self.presentation = None
                print("✅ Presentation closed")
                return True
            return False
        except Exception as e:
            print(f"❌ Error: {e}")
            return False
    
    def get_current_slide(self) -> int:
        """Get current slide number"""
        try:
            if self.slide_show:
                return self.slide_show.View.Slide.SlideIndex
            return 0
        except:
            return 0

# Global controller instance
ppt_controller = PowerPointController()

# ---------------- SPEECH MANAGER ----------------
class SpeechManager:
    """Thread-safe text-to-speech manager with interruption support"""
    def __init__(self):
        self.lock = threading.Lock()
        self.current_engine = None

    def speak(self, text):
        """Clean TTS with interruption handling"""
        print(f"\n🔊 Assistant: {text}")

        with self.lock:
            is_speaking.set()
            speech_interrupted.clear()

            try:
                self.current_engine = pyttsx3.init('sapi5')
                self.current_engine.setProperty('rate', 190)
                self.current_engine.setProperty('volume', 0.9)
                self.current_engine.say(text)
                self.current_engine.startLoop(False)

                while is_speaking.is_set():
                    if speech_interrupted.is_set():
                        self.current_engine.stop()
                        print("⏸️ Speech interrupted")
                        break
                    self.current_engine.iterate()

                self.current_engine.endLoop()

            except Exception as e:
                print(f"[TTS Error] {e}")

            finally:
                try:
                    if self.current_engine:
                        self.current_engine.stop()
                except:
                    pass

                self.current_engine = None
                is_speaking.clear()
                speech_interrupted.clear()
    
    def interrupt(self):
        """Interrupt ongoing speech"""
        speech_interrupted.set()
        if self.current_engine:
            try:
                self.current_engine.stop()
            except:
                pass

speech_manager = SpeechManager()

# ---------------- BACKGROUND LISTENER ----------------
def start_listener():
    """Initialize continuous speech recognition in background"""
    recognizer = sr.Recognizer()
    recognizer.energy_threshold = 4000
    recognizer.dynamic_energy_threshold = True
    recognizer.pause_threshold = 1.0
    
    mic = sr.Microphone()
    
    with mic as source:
        print("🎤 Calibrating microphone...")
        recognizer.adjust_for_ambient_noise(source, duration=2)
        print("✅ Calibration complete")

    def callback(recognizer, audio):
        try:
            text = recognizer.recognize_google(audio).strip().lower()
            if not text:
                return
            
            print(f"👂 Heard: {text}")
            
            if is_speaking.is_set():
                print("🛑 Interrupting speech")
                speech_manager.interrupt()
                user_input_queue.put(text)
            else:
                if not is_agent_awake.is_set():
                    if any(wake_word in text for wake_word in WAKE_WORDS):
                        print("⏰ Wake word detected!")
                        user_input_queue.put(text)
                        is_agent_awake.set()
                else:
                    user_input_queue.put(text)
                
        except sr.UnknownValueError:
            pass
        except Exception as e:
            print(f"[Listener Error] {e}")

    recognizer.listen_in_background(mic, callback, phrase_time_limit=15)
    print("🎧 Background listener active\n")

# ---------------- STATE DEFINITION ----------------
class AgentState(TypedDict):
    user_input: str
    intent: str
    response: str
    should_continue: bool
    waiting_for_input: bool

# ---------------- GRAPH NODES ----------------
def wake_up_node(state: AgentState) -> AgentState:
    """Initial greeting when agent activates"""
    
    if not has_greeted.is_set():
        current_time = datetime.now(ZoneInfo(TIMEZONE))
        hour = current_time.hour
        
        if 5 <= hour < 12:
            greeting = "Good morning"
        elif 12 <= hour < 17:
            greeting = "Good afternoon"
        else:
            greeting = "Good evening"
        
        prompt = f"""You're a friendly AI assistant being activated for the first time. 
        Current time: {current_time.strftime('%I:%M %p')}
        Give a brief, warm {greeting} greeting (1-2 sentences) and ask how you can help."""
        
        response = model_call(prompt)
        state['response'] = response
        
        conversation_history.append({"role": "assistant", "content": response})
        speech_manager.speak(response)
        
        has_greeted.set()
    
    state['should_continue'] = True
    state['waiting_for_input'] = True
    return state

def listen_node(state: AgentState) -> AgentState:
    """Wait for user input"""
    print("\n⏳ Waiting for your input...")
    
    try:
        user_input = user_input_queue.get_nowait()
        state['user_input'] = user_input
        state['waiting_for_input'] = False
        print(f"✅ Got queued input: {user_input}")
        return state
    except queue.Empty:
        pass
    
    timeout_counter = 0
    max_timeout = 120
    
    while True:
        try:
            user_input = user_input_queue.get(timeout=1)
            state['user_input'] = user_input
            state['waiting_for_input'] = False
            print(f"✅ Got input: {user_input}")
            break
            
        except queue.Empty:
            timeout_counter += 1
            
            if timeout_counter >= max_timeout:
                prompt_msg = "Are you still there?"
                speech_manager.speak(prompt_msg)
                timeout_counter = 0
            
            if shutdown_flag.is_set():
                state['should_continue'] = False
                state['waiting_for_input'] = False
                break
            
            continue
    
    return state

def classify_intent_node(state: AgentState) -> AgentState:
    """Classify user intent"""
    user_input = state['user_input']
    
    context = "\n".join([
        f"{msg['role'].title()}: {msg['content']}" 
        for msg in conversation_history[-4:]
    ]) if conversation_history else "No prior conversation"
    
    prompt = f"""Classify the following input into ONE category:

Categories:
- 'chat': General conversation, questions, or requests
- 'presentation': Requests to present/show slides, navigate (next, previous, open), OR questions about document content (what, explain, describe)
- 'sleep': Wants to end session (goodbye, that's all, sleep, stop, exit)

Recent context:
{context}

Current input: "{user_input}"

Respond with ONLY ONE WORD: chat, presentation, or sleep"""
    
    intent = model_call(prompt).lower().strip()
    
    if intent not in ['chat', 'presentation', 'sleep']:
        intent = 'chat'
    
    state['intent'] = intent
    print(f"🎯 Intent: {intent}")
    
    return state

def chat_node(state: AgentState) -> AgentState:
    """Handle general conversation"""
    user_input = state['user_input']
    
    context = "\n".join([
        f"{msg['role'].title()}: {msg['content']}" 
        for msg in conversation_history[-8:]
    ]) if conversation_history else "No prior conversation"
    
    prompt = f"""You are a helpful, conversational AI assistant.

Recent conversation:
{context}

User: {user_input}

Provide a natural, concise response (2-3 sentences max). Be helpful and friendly."""
    
    response = model_call(prompt)
    state['response'] = response
    
    conversation_history.append({"role": "user", "content": user_input})
    conversation_history.append({"role": "assistant", "content": response})
    
    speech_manager.speak(response)
    
    state['should_continue'] = True
    state['waiting_for_input'] = True
    return state

def presentation_node(state: AgentState) -> AgentState:
    """Handle presentation control AND document Q&A using RAG"""
    global current_document, current_page
    
    user_input = state['user_input']
    user_input_lower = user_input.lower()
    
    response = ""
    
    try:
        # ========== FUNCTIONALITY 1: DOCUMENT INDEXING ==========
        if 'index' in user_input_lower and ('document' in user_input_lower or 'presentation' in user_input_lower):
            if not RAG_AVAILABLE or not rag_manager:
                response = "RAG features are not available. Please install: pip install PyPDF2 python-pptx chromadb sentence-transformers"
            else:
                filename = "NeuralTwin.pptx"
                result = rag_manager.ingest_document(filename)
                
                if result['success']:
                    if result.get('message') == "Already indexed":
                        response = f"Failed to index document: {result.get('error', 'Unknown error')}"
        
        # ========== FUNCTIONALITY 2: RAG Q&A (Questions about content) ==========
        elif any(word in user_input_lower for word in ['what', 'explain', 'tell me about', 'describe', 'how', 'why', 'who', 'when', 'where']) and not any(word in user_input_lower for word in ['open', 'next', 'previous', 'close', 'slide', 'page']):
            if not RAG_AVAILABLE or not rag_manager:
                response = "I can't answer questions about document content without RAG features installed."
            elif not rag_manager.get_indexed_documents():
                response = "No documents are indexed yet. Please say 'index document' first, then I can answer your questions about the content."
            else:
                # Use current document or first indexed document
                filename = current_document if current_document and rag_manager.is_indexed(current_document) else rag_manager.get_indexed_documents()[0]
                
                # Retrieve relevant context
                search_results = rag_manager.retrieve_context(user_input, filename=filename, top_k=3)
                
                if search_results:
                    # Build context from retrieved chunks
                    context_text = "\n\n".join([
                        f"[Page/Slide {r['page_num']}]:\n{r['text'][:500]}" 
                        for r in search_results
                    ])
                    
                    # Generate answer using LLM with retrieved context
                    prompt = f"""Based on the following content from the presentation "{filename}", answer the user's question naturally and conversationally.

Retrieved context from the document:
{context_text}

User question: {user_input}

Provide a clear, conversational answer (2-4 sentences). Reference specific pages/slides if helpful. If the context doesn't fully answer the question, say what you can based on the available information."""
                    
                    response = model_call(prompt)
                else:
                    response = f"I couldn't find relevant information in {filename} to answer that question. Try rephrasing or ask about a different topic."
        
        # ========== FUNCTIONALITY 3: PRESENTATION CONTROL ==========
        
        # Command: Open presentation
        elif any(word in user_input_lower for word in ['open', 'start', 'show', 'present', 'load']):
            filename = "NeuralTwin.pptx"
            
            if ppt_controller.open_presentation(filename):
                current_document = filename
                current_page = 1
                response = f"Opening presentation: {filename}"
                
                # Auto-index if RAG available and not already indexed
                if RAG_AVAILABLE and rag_manager and not rag_manager.is_indexed(filename):
                    print("🔄 Auto-indexing document for Q&A...")
                    result = rag_manager.ingest_document(filename)
                    if result['success']:
                        response += f" I've also indexed it, so you can ask me questions about the content."
            else:
                response = f"I couldn't open {filename}. Please check if the file exists in {DOCUMENTS_FOLDER}"
        
        # Command: Next slide
        elif any(word in user_input_lower for word in ['next', 'forward', 'next slide']):
            if ppt_controller.slide_show:
                if ppt_controller.next_slide():
                    current_page = ppt_controller.get_current_slide()
                    response = f"Next slide: {current_page}"
                else:
                    response = "We're on the last slide"
            else:
                response = "No presentation is open. Please ask me to open a presentation first."
        
        # Command: Previous slide
        elif any(word in user_input_lower for word in ['previous', 'back', 'previous slide', 'go back']):
            if ppt_controller.slide_show:
                if ppt_controller.previous_slide():
                    current_page = ppt_controller.get_current_slide()
                    response = f"Previous slide: {current_page}"
                else:
                    response = "We're on the first slide"
            else:
                response = "No presentation is open."
        
        # Command: Go to specific slide
        elif 'slide' in user_input_lower and any(char.isdigit() for char in user_input):
            if ppt_controller.slide_show:
                import re
                numbers = re.findall(r'\d+', user_input)
                if numbers:
                    slide_num = int(numbers[0])
                    if ppt_controller.goto_slide(slide_num):
                        current_page = slide_num
                        response = f"Going to slide {slide_num}"
                    else:
                        response = f"Slide {slide_num} doesn't exist"
                else:
                    response = "I couldn't understand which slide number"
            else:
                response = "No presentation is open."
        
        # Command: Close presentation
        elif any(word in user_input_lower for word in ['close', 'exit presentation', 'stop presenting']):
            if ppt_controller.close_presentation():
                response = "Closing presentation"
                current_document = None
                current_page = 0
            else:
                response = "No presentation to close"
        
        # Default: Show capabilities
        else:
            if RAG_AVAILABLE and rag_manager:
                response = "I can help with presentations! Try: 'open presentation' to display slides, 'next slide' to navigate, 'index document' to enable Q&A, or ask me questions about the content after indexing."
            else:
                response = "I can help with presentations! Try: 'open presentation', 'next slide', 'previous slide', or 'close presentation'."
        
    except Exception as e:
        print(f"[Presentation Error] {e}")
        traceback.print_exc()
        response = "I encountered an error with presentation control."
    
    state['response'] = response
    
    conversation_history.append({"role": "user", "content": user_input})
    conversation_history.append({"role": "assistant", "content": response})
    
    speech_manager.speak(response)
    
    state['should_continue'] = True
    state['waiting_for_input'] = True
    return state

def sleep_node(state: AgentState) -> AgentState:
    """Handle session end"""
    global current_document, current_page
    
    current_time = datetime.now(ZoneInfo(TIMEZONE))
    hour = current_time.hour
    
    if 5 <= hour < 17:
        farewell = "Have a great day"
    else:
        farewell = "Have a wonderful evening"
    
    prompt = f"""The user is ending the session. 
    Give a brief, warm farewell (1 sentence). 
    Use this context: {farewell}. 
    Add a caring note."""
    
    response = model_call(prompt)
    state['response'] = response
    
    conversation_history.append({"role": "user", "content": state['user_input']})
    conversation_history.append({"role": "assistant", "content": response})
    speech_manager.speak(response)
    
    # Close presentation if open
    if ppt_controller.slide_show:
        ppt_controller.close_presentation()
    
    current_document = None
    current_page = 0
    
    is_agent_awake.clear()
    has_greeted.clear()
    
    state['should_continue'] = False
    state['waiting_for_input'] = False
    print("💤 Agent going to sleep. Say wake word to reactivate.\n")
    
    return state

# ---------------- ROUTING LOGIC ----------------
def route_by_intent(state: AgentState) -> str:
    intent = state.get('intent', 'chat')
    route_map = {'chat': 'chat', 'presentation': 'presentation', 'sleep': 'sleep'}
    return route_map.get(intent, 'chat')

def route_after_processing(state: AgentState) -> str:
    return 'continue' if state.get('should_continue', False) else 'end'

# ---------------- BUILD GRAPH ----------------
def build_graph():
    workflow = StateGraph(AgentState)
    
    workflow.add_node("wake_up", wake_up_node)
    workflow.add_node("listen", listen_node)
    workflow.add_node("classify", classify_intent_node)
    workflow.add_node("chat", chat_node)
    workflow.add_node("presentation", presentation_node)
    workflow.add_node("sleep", sleep_node)
    
    workflow.set_entry_point("wake_up")
    workflow.add_edge("wake_up", "listen")
    workflow.add_edge("listen", "classify")
    
    workflow.add_conditional_edges("classify", route_by_intent,
        {"chat": "chat", "presentation": "presentation", "sleep": "sleep"})
    
    workflow.add_conditional_edges("chat", route_after_processing,
        {"continue": "listen", "end": END})
    
    workflow.add_conditional_edges("presentation", route_after_processing,
        {"continue": "listen", "end": END})
    
    workflow.add_edge("sleep", END)
    
    return workflow.compile()

# ---------------- MAIN ----------------
def main():
    print("=" * 60)
    print("✨ VOICE AGENT SYSTEM (Laptop Standalone)")
    print("=" * 60)
    print(f"📂 Documents folder: {DOCUMENTS_FOLDER}")
    print(f"💾 Vector DB: {VECTOR_DB_PATH}")
    print(f"💤 Say one of these to wake: {', '.join(WAKE_WORDS)}")
    print("🎤 Speak clearly - I'll wait for you")
    print("🛑 Interrupt me anytime while speaking")
    print("💬 Say 'goodbye' to sleep")
    print()
    
    if RAG_AVAILABLE and rag_manager:
        print("✅ RAG system ready for document Q&A")
    else:
        print("⚠️ RAG features disabled (missing dependencies)")
    
    print("⌨️  Press Ctrl+C to exit\n")
    print("=" * 60)
    print()
    
    if not DOCUMENTS_FOLDER.exists():
        print(f"⚠️ Warning: {DOCUMENTS_FOLDER} not found!")
        print("Please update DOCUMENTS_FOLDER in the script\n")
    
    start_listener()
    agent_graph = build_graph()
    
    while not shutdown_flag.is_set():
        try:
            if not is_agent_awake.is_set():
                time.sleep(0.5)
                continue
            
            state = AgentState(
                user_input="", intent="", response="",
                should_continue=True, waiting_for_input=False
            )
            
            try:
                result = agent_graph.invoke(state)
                
                if not result.get('should_continue', False):
                    print("\n" + "=" * 60)
                    print("💤 Conversation ended. Listening for wake word...")
                    print("=" * 60 + "\n")
                    
            except Exception as e:
                print(f"[Graph Error] {e}")
                traceback.print_exc()
                time.sleep(1)
                
        except KeyboardInterrupt:
            print("\n\n👋 Shutting down...")
            if is_agent_awake.is_set():
                speech_manager.speak("Goodbye!")
            break
        except Exception as e:
            print(f"[Main Error] {e}")
            traceback.print_exc()
            time.sleep(1)
    
    shutdown_flag.set()
    print("✨ System terminated.")

if __name__ == "__main__":
    import sys
    if sys.platform != "win32":
        print("❌ This script requires Windows (pywin32)")
        sys.exit(1)
    
    main()
